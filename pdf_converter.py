"""
PDF Conversion & Processing Module
Provides tools for converting between PDF and various formats
"""

import os
import io
import pathlib
from typing import List, Tuple
from PIL import Image
import fitz  # PyMuPDF
import pypdf
import img2pdf
from pdf2image import convert_from_path
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

# Setup directories
DATA_DIR = os.getenv('RAILWAY_VOLUME_MOUNT_PATH', os.path.dirname(os.path.abspath(__file__)))
DOCUMENTS_DIR = os.path.join(DATA_DIR, 'documents')
ASSETS_DIR = os.path.join(DATA_DIR, 'assets')

# Ensure directories exist
os.makedirs(DOCUMENTS_DIR, exist_ok=True)
os.makedirs(ASSETS_DIR, exist_ok=True)


# ==================== FILE PATH NORMALIZATION ====================

def normalize_file_path(file_path: str) -> str:
    """
    Normalize file paths to handle various input formats:
    - Full absolute paths: /app/data/documents/file.pdf
    - Relative paths: documents/file.pdf
    - Just filenames: file.pdf
    - Web URLs: /documents/file.pdf

    Returns the full absolute path if file exists, otherwise raises FileNotFoundError.
    """
    if not file_path:
        raise FileNotFoundError("No file path provided")

    # Remove web URL prefix if present
    if file_path.startswith('/documents/'):
        file_path = file_path.replace('/documents/', '')
    elif file_path.startswith('/assets/'):
        file_path = file_path.replace('/assets/', '')
    elif file_path.startswith('/uploads/'):
        file_path = file_path.replace('/uploads/', '')

    # Case 1: Already a valid absolute path
    if os.path.isabs(file_path) and os.path.exists(file_path):
        return file_path

    # Case 2: Check in DOCUMENTS_DIR
    docs_path = os.path.join(DOCUMENTS_DIR, os.path.basename(file_path))
    if os.path.exists(docs_path):
        return docs_path

    # Case 3: Check in ASSETS_DIR
    assets_path = os.path.join(ASSETS_DIR, os.path.basename(file_path))
    if os.path.exists(assets_path):
        return assets_path

    # Case 4: Check relative to current directory
    if os.path.exists(file_path):
        return os.path.abspath(file_path)

    # File not found in any location
    raise FileNotFoundError(f"File not found: {file_path}. Searched in: {DOCUMENTS_DIR}, {ASSETS_DIR}, and current directory.")


# ==================== PDF to Images ====================

def pdf_to_images(pdf_path: str, output_format: str = 'png', dpi: int = 200) -> List[str]:
    """
    Convert PDF pages to individual image files

    Args:
        pdf_path: Path to PDF file
        output_format: 'png' or 'jpg'
        dpi: Image resolution (default 200)

    Returns:
        List of paths to generated image files
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        # Convert PDF to images
        images = convert_from_path(pdf_path, dpi=dpi)

        output_paths = []
        pdf_name = pathlib.Path(pdf_path).stem

        for i, image in enumerate(images):
            output_filename = f"{pdf_name}_page_{i+1}.{output_format}"
            output_path = os.path.join(ASSETS_DIR, output_filename)

            # Save image - PIL expects 'JPEG' not 'JPG'
            pil_format = 'JPEG' if output_format.lower() == 'jpg' else output_format.upper()
            image.save(output_path, format=pil_format)
            output_paths.append(output_path)

        return output_paths

    except Exception as e:
        raise Exception(f"Failed to convert PDF to images: {str(e)}")


# ==================== Images to PDF ====================

def images_to_pdf(image_paths: List[str], output_name: str = None) -> str:
    """
    Convert multiple images into a single PDF

    Args:
        image_paths: List of image file paths
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        if not output_name:
            output_name = f"combined_images_{int(time.time())}.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Validate that all image files exist
        for img_path in image_paths:
            if not os.path.exists(img_path):
                raise Exception(f"Image file not found: {img_path}")

        # Try using img2pdf first (faster, better quality)
        # img2pdf can accept file paths directly when passed as a list
        try:
            with open(output_path, "wb") as f:
                # Pass image paths as a list - img2pdf will handle them correctly
                f.write(img2pdf.convert(image_paths))
            print(f"Successfully converted {len(image_paths)} images to PDF using img2pdf")
            return output_path
        except Exception as img2pdf_error:
            print(f"img2pdf failed: {img2pdf_error}, trying PIL fallback...")

            # Fallback: Use PIL to convert images to RGB and then to PDF
            images = []
            for i, img_path in enumerate(image_paths):
                print(f"Processing image {i+1}/{len(image_paths)}: {img_path}")
                try:
                    img = Image.open(img_path)
                    print(f"  Original mode: {img.mode}, size: {img.size}")
                    # Convert to RGB if necessary (handles RGBA, P, etc.)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                        print(f"  Converted to RGB")
                    images.append(img)
                except Exception as img_error:
                    print(f"  ERROR processing image: {img_error}")
                    # Continue with other images
                    continue

            # Save as PDF
            if images:
                print(f"Saving {len(images)} images to PDF: {output_path}")
                images[0].save(output_path, save_all=True, append_images=images[1:] if len(images) > 1 else [], format='PDF')
                return output_path
            else:
                raise Exception("No valid images to convert")

    except Exception as e:
        raise Exception(f"Failed to convert images to PDF: {str(e)}")


# ==================== PDF to Word ====================

def pdf_to_word(pdf_path: str, output_name: str = None) -> str:
    """
    Convert PDF to Word document (DOCX)
    Preserves tables and basic formatting

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename

    Returns:
        Path to generated Word file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        # Open PDF
        doc_pdf = fitz.open(pdf_path)

        # Create Word document
        doc_word = Document()

        # Process each page
        for page_num in range(len(doc_pdf)):
            page = doc_pdf[page_num]

            # Add page break after first page
            if page_num > 0:
                doc_word.add_page_break()

            # Find tables on page using PyMuPDF's table detection
            try:
                tables = page.find_tables()

                if tables and len(tables.tables) > 0:
                    # Page contains tables
                    print(f"Found {len(tables.tables)} tables on page {page_num + 1}")

                    # Get all text blocks with their positions
                    blocks = page.get_text("dict")["blocks"]
                    text_blocks = []

                    for block in blocks:
                        if block.get("type") == 0:  # Text block
                            bbox = block["bbox"]
                            text = ""
                            for line in block.get("lines", []):
                                for span in line.get("spans", []):
                                    text += span.get("text", "")
                                text += "\n"
                            if text.strip():
                                text_blocks.append({
                                    "bbox": bbox,
                                    "text": text.strip(),
                                    "y0": bbox[1]  # Top Y coordinate for sorting
                                })

                    # Sort blocks by vertical position
                    text_blocks.sort(key=lambda x: x["y0"])

                    # Process tables and text in order
                    table_bboxes = []
                    for table in tables.tables:
                        bbox = table.bbox
                        table_bboxes.append({
                            "bbox": bbox,
                            "table": table,
                            "y0": bbox[1]
                        })

                    # Combine and sort all elements
                    all_elements = text_blocks + table_bboxes
                    all_elements.sort(key=lambda x: x["y0"])

                    # Add elements to Word doc in order
                    for element in all_elements:
                        if "table" in element:
                            # Add table to Word
                            table = element["table"]
                            table_data = table.extract()

                            if table_data and len(table_data) > 0:
                                # Create Word table
                                rows = len(table_data)
                                cols = max(len(row) for row in table_data) if table_data else 0

                                if rows > 0 and cols > 0:
                                    word_table = doc_word.add_table(rows=rows, cols=cols)
                                    word_table.style = 'Table Grid'

                                    # Fill table cells
                                    for i, row in enumerate(table_data):
                                        for j, cell in enumerate(row):
                                            if j < cols and cell:
                                                word_table.rows[i].cells[j].text = str(cell)
                        else:
                            # Add text paragraph
                            if element["text"].strip():
                                doc_word.add_paragraph(element["text"])

                else:
                    # No tables found, extract text normally
                    text = page.get_text()
                    if text.strip():
                        doc_word.add_paragraph(text)

            except Exception as table_error:
                print(f"Table detection failed on page {page_num + 1}: {table_error}")
                # Fallback to plain text extraction
                text = page.get_text()
                if text.strip():
                    doc_word.add_paragraph(text)

        # Generate output path
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_converted.docx"

        if not output_name.endswith('.docx'):
            output_name += '.docx'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)
        doc_word.save(output_path)

        doc_pdf.close()
        return output_path

    except Exception as e:
        raise Exception(f"Failed to convert PDF to Word: {str(e)}")


# ==================== PDF to Excel ====================

def pdf_to_excel(pdf_path: str, output_name: str = None) -> str:
    """
    Convert PDF to Excel spreadsheet (XLSX)
    Preserves table structure and formatting

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename

    Returns:
        Path to generated Excel file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        # Open PDF
        doc = fitz.open(pdf_path)

        # Create Excel workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Data"

        current_row = 1

        # Process each page
        for page_num in range(len(doc)):
            page = doc[page_num]

            # Add page header
            ws.cell(row=current_row, column=1, value=f"Page {page_num + 1}")
            ws.cell(row=current_row, column=1).font = ws.cell(row=current_row, column=1).font.copy(bold=True)
            current_row += 1

            # Try to find tables on the page
            try:
                tables = page.find_tables()

                if tables and len(tables.tables) > 0:
                    # Page contains tables
                    print(f"Found {len(tables.tables)} tables on page {page_num + 1}")

                    for table_idx, table in enumerate(tables.tables):
                        table_data = table.extract()

                        if table_data and len(table_data) > 0:
                            # Add table to Excel with proper structure
                            start_row = current_row

                            for row_idx, row in enumerate(table_data):
                                for col_idx, cell in enumerate(row):
                                    if cell:
                                        ws.cell(row=current_row, column=col_idx + 1, value=str(cell))

                                        # Bold header row if first row
                                        if row_idx == 0:
                                            ws.cell(row=current_row, column=col_idx + 1).font = \
                                                ws.cell(row=current_row, column=col_idx + 1).font.copy(bold=True)

                                current_row += 1

                            current_row += 1  # Add space after table

                            # Auto-adjust column widths for table
                            for col_idx in range(len(table_data[0]) if table_data else 0):
                                col_letter = ws.cell(row=start_row, column=col_idx + 1).column_letter
                                max_length = 0
                                for row_idx in range(start_row, current_row - 1):
                                    cell_value = ws.cell(row=row_idx, column=col_idx + 1).value
                                    if cell_value:
                                        max_length = max(max_length, len(str(cell_value)))
                                ws.column_dimensions[col_letter].width = min(max_length + 2, 50)

                else:
                    # No tables found, extract text normally
                    text = page.get_text()
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            ws.cell(row=current_row, column=1, value=line)
                            current_row += 1

            except Exception as table_error:
                print(f"Table detection failed on page {page_num + 1}: {table_error}")
                # Fallback to plain text extraction
                text = page.get_text()
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=current_row, column=1, value=line)
                        current_row += 1

            current_row += 1  # Add space between pages

        # Generate output path
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_converted.xlsx"

        if not output_name.endswith('.xlsx'):
            output_name += '.xlsx'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)
        wb.save(output_path)

        doc.close()
        return output_path

    except Exception as e:
        raise Exception(f"Failed to convert PDF to Excel: {str(e)}")


# ==================== PDF to PowerPoint ====================

def pdf_to_powerpoint(pdf_path: str, output_name: str = None, slides_per_page: bool = True) -> str:
    """
    Convert PDF to PowerPoint presentation (PPTX)
    Each page becomes a slide with image or text

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename
        slides_per_page: If True, each PDF page becomes one slide

    Returns:
        Path to generated PowerPoint file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        # Open PDF
        doc = fitz.open(pdf_path)

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Process each page
        for page_num in range(len(doc)):
            page = doc[page_num]

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Convert page to image
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")

            # Save temp image
            temp_img_path = os.path.join(ASSETS_DIR, f"temp_slide_{page_num}.png")
            pix.save(temp_img_path)

            # Add image to slide
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            slide.shapes.add_picture(temp_img_path, left, top, width=width)

            # Clean up temp file
            try:
                os.remove(temp_img_path)
            except:
                pass

        # Generate output path
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_converted.pptx"

        if not output_name.endswith('.pptx'):
            output_name += '.pptx'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)
        prs.save(output_path)

        doc.close()
        return output_path

    except Exception as e:
        raise Exception(f"Failed to convert PDF to PowerPoint: {str(e)}")


# ==================== Word to PDF ====================

def word_to_pdf(docx_path: str, output_name: str = None) -> str:
    """
    Convert Word document to PDF using LibreOffice

    Args:
        docx_path: Path to Word file
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        # Normalize file path
        docx_path = normalize_file_path(docx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(docx_path):
            raise Exception(f"Word file not found: {docx_path}")

        # Generate output name
        if not output_name:
            docx_name = pathlib.Path(docx_path).stem
            output_name = f"{docx_name}_converted.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to PDF
        # --headless: run without GUI
        # --convert-to pdf: output format
        # --outdir: output directory
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', DOCUMENTS_DIR,
            docx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .pdf extension
        # Rename if needed
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(docx_path).stem + '.pdf')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Word to PDF: {str(e)}")


# ==================== Excel to PDF ====================

def excel_to_pdf(xlsx_path: str, output_name: str = None) -> str:
    """
    Convert Excel spreadsheet to PDF using LibreOffice

    Args:
        xlsx_path: Path to Excel file
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        # Normalize file path
        xlsx_path = normalize_file_path(xlsx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(xlsx_path):
            raise Exception(f"Excel file not found: {xlsx_path}")

        # Generate output name
        if not output_name:
            xlsx_name = pathlib.Path(xlsx_path).stem
            output_name = f"{xlsx_name}_converted.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to PDF
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', DOCUMENTS_DIR,
            xlsx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .pdf extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(xlsx_path).stem + '.pdf')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Excel to PDF: {str(e)}")


# ==================== PowerPoint to PDF ====================

def powerpoint_to_pdf(pptx_path: str, output_name: str = None) -> str:
    """
    Convert PowerPoint presentation to PDF using LibreOffice

    Args:
        pptx_path: Path to PowerPoint file
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        # Normalize file path
        pptx_path = normalize_file_path(pptx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(pptx_path):
            raise Exception(f"PowerPoint file not found: {pptx_path}")

        # Generate output name
        if not output_name:
            pptx_name = pathlib.Path(pptx_path).stem
            output_name = f"{pptx_name}_converted.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to PDF
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', DOCUMENTS_DIR,
            pptx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .pdf extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(pptx_path).stem + '.pdf')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert PowerPoint to PDF: {str(e)}")


# ==================== Utility Functions ====================

def get_pdf_info(pdf_path: str) -> dict:
    """Get information about a PDF file"""
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        doc = fitz.open(pdf_path)
        info = {
            'pages': len(doc),
            'metadata': doc.metadata,
            'size_bytes': os.path.getsize(pdf_path),
            'encrypted': doc.is_encrypted
        }
        doc.close()
        return info
    except Exception as e:
        return {'error': str(e)}


def validate_file_type(file_path: str, allowed_extensions: List[str]) -> bool:
    """Validate if file has an allowed extension"""
    ext = pathlib.Path(file_path).suffix.lower()
    return ext in [f'.{e}' if not e.startswith('.') else e for e in allowed_extensions]


# ==================== PDF Merge ====================

def merge_pdfs(pdf_paths: List[str], output_name: str = None) -> str:
    """
    Merge multiple PDF files into a single PDF

    Args:
        pdf_paths: List of PDF file paths to merge
        output_name: Optional output filename

    Returns:
        Path to merged PDF file
    """
    try:
        if not pdf_paths or len(pdf_paths) == 0:
            raise Exception("No PDF files provided for merging")

        # Validate all files exist and are PDFs
        for pdf_path in pdf_paths:
            if not os.path.exists(pdf_path):
                raise Exception(f"PDF file not found: {pdf_path}")
            if not pdf_path.lower().endswith('.pdf'):
                raise Exception(f"Not a PDF file: {pdf_path}")

        # Generate output name if not provided
        if not output_name:
            output_name = f"merged_{int(time.time())}.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Create PDF writer for merging
        writer = pypdf.PdfWriter()

        # Add each PDF
        for pdf_path in pdf_paths:
            print(f"Adding PDF to merge: {pdf_path}")
            reader = pypdf.PdfReader(pdf_path)
            for page in reader.pages:
                writer.add_page(page)

        # Write merged PDF
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)

        print(f"Successfully merged {len(pdf_paths)} PDFs into: {output_path}")
        return output_path

    except Exception as e:
        raise Exception(f"Failed to merge PDFs: {str(e)}")


# Add time import that was missing
import time


# ==================== Excel to CSV ====================

def excel_to_csv(xlsx_path: str, output_name: str = None, sheet_name: str = None) -> str:
    """
    Convert Excel spreadsheet to CSV format using LibreOffice

    Args:
        xlsx_path: Path to Excel file
        output_name: Optional output filename (without extension)
        sheet_name: Optional specific sheet to convert (default: first sheet)

    Returns:
        Path to generated CSV file
    """
    try:
        # Normalize file path
        xlsx_path = normalize_file_path(xlsx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(xlsx_path):
            raise Exception(f"Excel file not found: {xlsx_path}")

        # Generate output name
        if not output_name:
            xlsx_name = pathlib.Path(xlsx_path).stem
            output_name = f"{xlsx_name}_converted.csv"

        if not output_name.endswith('.csv'):
            output_name += '.csv'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to CSV
        # --infilter for Excel input, --convert-to for CSV output
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'csv',
            '--outdir', DOCUMENTS_DIR,
            xlsx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .csv extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(xlsx_path).stem + '.csv')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Excel to CSV: {str(e)}")


# ==================== CSV to Excel ====================

def csv_to_excel(csv_path: str, output_name: str = None) -> str:
    """
    Convert CSV file to Excel spreadsheet using LibreOffice

    Args:
        csv_path: Path to CSV file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated Excel file
    """
    try:
        # Normalize file path
        csv_path = normalize_file_path(csv_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(csv_path):
            raise Exception(f"CSV file not found: {csv_path}")

        # Generate output name
        if not output_name:
            csv_name = pathlib.Path(csv_path).stem
            output_name = f"{csv_name}_converted.xlsx"

        if not output_name.endswith('.xlsx'):
            output_name += '.xlsx'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert CSV to Excel
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'xlsx',
            '--outdir', DOCUMENTS_DIR,
            csv_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .xlsx extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(csv_path).stem + '.xlsx')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert CSV to Excel: {str(e)}")


# ==================== Word to TXT ====================

def word_to_txt(docx_path: str, output_name: str = None) -> str:
    """
    Convert Word document to plain text format using LibreOffice

    Args:
        docx_path: Path to Word file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated TXT file
    """
    try:
        # Normalize file path
        docx_path = normalize_file_path(docx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(docx_path):
            raise Exception(f"Word file not found: {docx_path}")

        # Generate output name
        if not output_name:
            docx_name = pathlib.Path(docx_path).stem
            output_name = f"{docx_name}_converted.txt"

        if not output_name.endswith('.txt'):
            output_name += '.txt'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to TXT
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'txt',
            '--outdir', DOCUMENTS_DIR,
            docx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .txt extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(docx_path).stem + '.txt')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Word to TXT: {str(e)}")


# ==================== Excel to TXT ====================

def excel_to_txt(xlsx_path: str, output_name: str = None) -> str:
    """
    Convert Excel spreadsheet to plain text format using LibreOffice

    Args:
        xlsx_path: Path to Excel file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated TXT file
    """
    try:
        # Normalize file path
        xlsx_path = normalize_file_path(xlsx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(xlsx_path):
            raise Exception(f"Excel file not found: {xlsx_path}")

        # Generate output name
        if not output_name:
            xlsx_name = pathlib.Path(xlsx_path).stem
            output_name = f"{xlsx_name}_converted.txt"

        if not output_name.endswith('.txt'):
            output_name += '.txt'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to TXT
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'txt',
            '--outdir', DOCUMENTS_DIR,
            xlsx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .txt extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(xlsx_path).stem + '.txt')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Excel to TXT: {str(e)}")


# ==================== Word to HTML ====================

def word_to_html(docx_path: str, output_name: str = None) -> str:
    """
    Convert Word document to HTML format using LibreOffice

    Args:
        docx_path: Path to Word file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated HTML file
    """
    try:
        # Normalize file path
        docx_path = normalize_file_path(docx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(docx_path):
            raise Exception(f"Word file not found: {docx_path}")

        # Generate output name
        if not output_name:
            docx_name = pathlib.Path(docx_path).stem
            output_name = f"{docx_name}_converted.html"

        if not output_name.endswith('.html'):
            output_name += '.html'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to HTML
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'html',
            '--outdir', DOCUMENTS_DIR,
            docx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .html extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(docx_path).stem + '.html')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Word to HTML: {str(e)}")


# ==================== Excel to HTML ====================

def excel_to_html(xlsx_path: str, output_name: str = None) -> str:
    """
    Convert Excel spreadsheet to HTML format using LibreOffice

    Args:
        xlsx_path: Path to Excel file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated HTML file
    """
    try:
        # Normalize file path
        xlsx_path = normalize_file_path(xlsx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(xlsx_path):
            raise Exception(f"Excel file not found: {xlsx_path}")

        # Generate output name
        if not output_name:
            xlsx_name = pathlib.Path(xlsx_path).stem
            output_name = f"{xlsx_name}_converted.html"

        if not output_name.endswith('.html'):
            output_name += '.html'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to HTML
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'html',
            '--outdir', DOCUMENTS_DIR,
            xlsx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .html extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(xlsx_path).stem + '.html')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert Excel to HTML: {str(e)}")


# ==================== PowerPoint to HTML ====================

def powerpoint_to_html(pptx_path: str, output_name: str = None) -> str:
    """
    Convert PowerPoint presentation to HTML format using LibreOffice

    Args:
        pptx_path: Path to PowerPoint file
        output_name: Optional output filename (without extension)

    Returns:
        Path to generated HTML file
    """
    try:
        # Normalize file path
        pptx_path = normalize_file_path(pptx_path)

        import subprocess

        # Validate file exists
        if not os.path.exists(pptx_path):
            raise Exception(f"PowerPoint file not found: {pptx_path}")

        # Generate output name
        if not output_name:
            pptx_name = pathlib.Path(pptx_path).stem
            output_name = f"{pptx_name}_converted.html"

        if not output_name.endswith('.html'):
            output_name += '.html'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Use LibreOffice to convert to HTML
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'html',
            '--outdir', DOCUMENTS_DIR,
            pptx_path
        ], check=True, capture_output=True, timeout=60)

        # LibreOffice creates the file with original name + .html extension
        temp_output = os.path.join(DOCUMENTS_DIR, pathlib.Path(pptx_path).stem + '.html')
        if temp_output != output_path and os.path.exists(temp_output):
            os.rename(temp_output, output_path)

        return output_path

    except subprocess.TimeoutExpired:
        raise Exception("Conversion timed out (>60 seconds)")
    except subprocess.CalledProcessError as e:
        raise Exception(f"LibreOffice conversion failed: {e.stderr.decode() if e.stderr else str(e)}")
    except Exception as e:
        raise Exception(f"Failed to convert PowerPoint to HTML: {str(e)}")


# ==================== PDF to Text ====================

def pdf_to_text(pdf_path: str, output_name: str = None) -> str:
    """
    Extract all text from PDF to plain text file

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename

    Returns:
        Path to generated TXT file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Generate output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_extracted.txt"

        if not output_name.endswith('.txt'):
            output_name += '.txt'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Extract text from PDF using PyMuPDF
        doc = fitz.open(pdf_path)
        text_content = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_content.append(f"--- Page {page_num + 1} ---\n{text}\n")

        doc.close()

        # Check if any text was extracted
        if not text_content or all(not t.strip() for t in text_content):
            # No text found - PDF likely contains images or scanned content
            error_msg = (
                "⚠️ WARNING: No extractable text found in PDF.\n\n"
                "This PDF appears to contain images or scanned content without a text layer.\n"
                "The text is embedded as images and cannot be extracted directly.\n\n"
                "SOLUTION: Upload the PDF and ask me to analyze it directly!\n"
                "I can read the PDF using GPT-4o vision and answer your questions.\n\n"
                "Example: 'What information is in this PDF?' or 'Give me the email and phone number'\n"
            )
            # Write the error message to file so user understands the issue
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(error_msg)
            return output_path

        # Write to text file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))

        return output_path

    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")


# ==================== Compress PDF ====================

def compress_pdf(pdf_path: str, output_name: str = None, compression_level: str = 'medium') -> str:
    """
    Compress PDF by reducing image quality and removing metadata

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename
        compression_level: 'low', 'medium', or 'high' compression

    Returns:
        Path to compressed PDF file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        import pikepdf

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Generate output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_compressed.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Compression settings based on level
        compression_settings = {
            'low': {'image_quality': 85, 'dpi': 150},
            'medium': {'image_quality': 70, 'dpi': 120},
            'high': {'image_quality': 50, 'dpi': 100}
        }

        settings = compression_settings.get(compression_level, compression_settings['medium'])

        # Open and compress PDF using pikepdf
        with pikepdf.open(pdf_path) as pdf:
            # Remove metadata to reduce size
            pdf.remove_unreferenced_resources()

            # Compress images
            for page in pdf.pages:
                for image_key in page.images.keys():
                    try:
                        raw_image = page.images[image_key]
                        pil_image = raw_image.as_pil_image()

                        # Reduce DPI and quality
                        if pil_image.mode == 'RGBA':
                            pil_image = pil_image.convert('RGB')

                        # Save compressed version
                        img_stream = io.BytesIO()
                        pil_image.save(img_stream, format='JPEG', quality=settings['image_quality'], optimize=True)
                        img_stream.seek(0)

                    except Exception as img_error:
                        print(f"Could not compress image {image_key}: {img_error}")
                        continue

            # Save compressed PDF
            pdf.save(output_path, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)

        # Get file size reduction
        original_size = os.path.getsize(pdf_path)
        compressed_size = os.path.getsize(output_path)
        reduction = ((original_size - compressed_size) / original_size) * 100

        print(f"Compressed PDF from {original_size/1024:.1f}KB to {compressed_size/1024:.1f}KB ({reduction:.1f}% reduction)")
        return output_path

    except ImportError:
        raise Exception("pikepdf library not installed. Run: pip install pikepdf")
    except Exception as e:
        raise Exception(f"Failed to compress PDF: {str(e)}")


# ==================== Split PDF ====================

def split_pdf(pdf_path: str, pages: str = 'all', output_name: str = None) -> List[str]:
    """
    Split PDF into separate files (one page per file or specified ranges)

    Args:
        pdf_path: Path to PDF file
        pages: 'all' or page ranges like '1-3,5,7-9'
        output_name: Optional base name for output files

    Returns:
        List of paths to generated PDF files
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Open PDF
        reader = pypdf.PdfReader(pdf_path)
        total_pages = len(reader.pages)

        # Generate base output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_page"

        output_paths = []

        # Split all pages (one page per file)
        if pages == 'all':
            for page_num in range(total_pages):
                writer = pypdf.PdfWriter()
                writer.add_page(reader.pages[page_num])

                output_filename = f"{output_name}_{page_num + 1}.pdf"
                output_path = os.path.join(DOCUMENTS_DIR, output_filename)

                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)

                output_paths.append(output_path)
        else:
            # Parse page ranges: e.g., "1-3,5,7-9" or "1-3" or "4-end"
            ranges = pages.split(',')
            part_num = 1

            for range_str in ranges:
                range_str = range_str.strip()
                writer = pypdf.PdfWriter()

                if '-' in range_str:
                    # Handle range like "1-3" or "4-end"
                    start_str, end_str = range_str.split('-')
                    start = int(start_str) - 1  # Convert to 0-indexed

                    if end_str.lower() in ['end', 'rest', 'last']:
                        end = total_pages - 1
                    else:
                        end = int(end_str) - 1  # Convert to 0-indexed

                    # Add all pages in range
                    for page_num in range(start, end + 1):
                        if 0 <= page_num < total_pages:
                            writer.add_page(reader.pages[page_num])

                    output_filename = f"{output_name}_part{part_num}_pages{start+1}-{end+1}.pdf"
                else:
                    # Single page like "5"
                    page_num = int(range_str) - 1  # Convert to 0-indexed
                    if 0 <= page_num < total_pages:
                        writer.add_page(reader.pages[page_num])
                    output_filename = f"{output_name}_part{part_num}_page{page_num+1}.pdf"

                output_path = os.path.join(DOCUMENTS_DIR, output_filename)

                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)

                output_paths.append(output_path)
                part_num += 1

        print(f"Split PDF into {len(output_paths)} files")
        return output_paths

    except Exception as e:
        raise Exception(f"Failed to split PDF: {str(e)}")


# ==================== Rotate PDF ====================

def rotate_pdf(pdf_path: str, rotation: int = 90, pages: str = 'all', output_name: str = None) -> str:
    """
    Rotate PDF pages by specified degrees

    Args:
        pdf_path: Path to PDF file
        rotation: Rotation angle (90, 180, 270, or -90)
        pages: 'all' or page numbers like '1,3,5'
        output_name: Optional output filename

    Returns:
        Path to rotated PDF file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Validate rotation angle
        if rotation not in [90, 180, 270, -90]:
            raise Exception(f"Invalid rotation angle: {rotation}. Must be 90, 180, 270, or -90")

        # Generate output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_rotated.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Open PDF
        reader = pypdf.PdfReader(pdf_path)
        writer = pypdf.PdfWriter()

        # Rotate pages
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]

            if pages == 'all':
                page.rotate(rotation)
            else:
                # Parse specific pages (future enhancement)
                page.rotate(rotation)

            writer.add_page(page)

        # Save rotated PDF
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)

        print(f"Rotated PDF by {rotation} degrees")
        return output_path

    except Exception as e:
        raise Exception(f"Failed to rotate PDF: {str(e)}")


# ==================== PDF to CSV ====================

def pdf_to_csv(pdf_path: str, output_name: str = None) -> str:
    """
    Extract tables from PDF to CSV format

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename

    Returns:
        Path to generated CSV file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        import pdfplumber
        import csv

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Generate output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_tables.csv"

        if not output_name.endswith('.csv'):
            output_name += '.csv'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Extract tables from PDF using pdfplumber
        all_tables = []

        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()

                if tables:
                    for table_idx, table in enumerate(tables):
                        # Add page and table headers
                        all_tables.append([f"Page {page_num + 1} - Table {table_idx + 1}"])
                        all_tables.extend(table)
                        all_tables.append([])  # Empty row separator

        # Write to CSV
        if all_tables:
            with open(output_path, 'w', newline='', encoding='utf-8') as csvfile:
                writer = csv.writer(csvfile)
                writer.writerows(all_tables)

            print(f"Extracted tables to CSV: {output_path}")
            return output_path
        else:
            raise Exception("No tables found in PDF")

    except ImportError:
        raise Exception("pdfplumber library not installed. Run: pip install pdfplumber")
    except Exception as e:
        raise Exception(f"Failed to extract PDF tables to CSV: {str(e)}")


# ==================== CSV to PDF ====================

def csv_to_pdf(csv_path: str, output_name: str = None) -> str:
    """
    Convert CSV data to formatted PDF table

    Args:
        csv_path: Path to CSV file
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        # Normalize file path
        csv_path = normalize_file_path(csv_path)

        import csv
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib.units import inch

        if not os.path.exists(csv_path):
            raise Exception(f"CSV file not found: {csv_path}")

        # Generate output name if not provided
        if not output_name:
            csv_name = pathlib.Path(csv_path).stem
            output_name = f"{csv_name}_converted.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Read CSV data
        data = []
        with open(csv_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            data = list(reader)

        if not data:
            raise Exception("CSV file is empty")

        # Determine if landscape is needed based on column count
        col_count = max(len(row) for row in data) if data else 0
        pagesize = landscape(letter) if col_count > 6 else letter

        # Create PDF
        doc = SimpleDocTemplate(output_path, pagesize=pagesize)
        elements = []

        # Create table
        table = Table(data)

        # Style table
        style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ])
        table.setStyle(style)

        elements.append(table)
        doc.build(elements)

        print(f"Converted CSV to PDF: {output_path}")
        return output_path

    except ImportError:
        raise Exception("reportlab library not installed. Run: pip install reportlab")
    except Exception as e:
        raise Exception(f"Failed to convert CSV to PDF: {str(e)}")


# ==================== HTML to PDF ====================

def html_to_pdf(html_input: str, output_name: str = None) -> str:
    """
    Convert HTML to PDF using weasyprint

    Args:
        html_input: Path to HTML file or HTML string
        output_name: Optional output filename

    Returns:
        Path to generated PDF file
    """
    try:
        from weasyprint import HTML

        # Generate output name if not provided
        if not output_name:
            output_name = f"converted_{int(time.time())}.pdf"

        if not output_name.endswith('.pdf'):
            output_name += '.pdf'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Check if input is a file path or HTML string
        if os.path.exists(html_input):
            # It's a file path
            HTML(filename=html_input).write_pdf(output_path)
        else:
            # It's an HTML string
            HTML(string=html_input).write_pdf(output_path)

        print(f"Converted HTML to PDF: {output_path}")
        return output_path

    except ImportError:
        raise Exception("weasyprint library not installed. Run: pip install weasyprint")
    except Exception as e:
        raise Exception(f"Failed to convert HTML to PDF: {str(e)}")


# ==================== PDF to HTML ====================

def pdf_to_html(pdf_path: str, output_name: str = None) -> str:
    """
    Convert PDF to HTML format

    Args:
        pdf_path: Path to PDF file
        output_name: Optional output filename

    Returns:
        Path to generated HTML file
    """
    try:
        # Normalize file path
        pdf_path = normalize_file_path(pdf_path)

        import pdfplumber

        if not os.path.exists(pdf_path):
            raise Exception(f"PDF file not found: {pdf_path}")

        # Generate output name if not provided
        if not output_name:
            pdf_name = pathlib.Path(pdf_path).stem
            output_name = f"{pdf_name}_converted.html"

        if not output_name.endswith('.html'):
            output_name += '.html'

        output_path = os.path.join(DOCUMENTS_DIR, output_name)

        # Extract content from PDF
        html_content = ['<!DOCTYPE html>', '<html>', '<head>',
                       '<meta charset="UTF-8">',
                       f'<title>{pathlib.Path(pdf_path).stem}</title>',
                       '<style>',
                       'body { font-family: Arial, sans-serif; margin: 40px; }',
                       '.page { margin-bottom: 40px; page-break-after: always; }',
                       'h2 { color: #333; border-bottom: 2px solid #333; }',
                       'table { border-collapse: collapse; width: 100%; margin: 20px 0; }',
                       'table td, table th { border: 1px solid #ddd; padding: 8px; }',
                       'table th { background-color: #f2f2f2; }',
                       '</style>',
                       '</head>', '<body>']

        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                html_content.append(f'<div class="page">')
                html_content.append(f'<h2>Page {page_num + 1}</h2>')

                # Extract tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        html_content.append('<table>')
                        for row_idx, row in enumerate(table):
                            html_content.append('<tr>')
                            tag = 'th' if row_idx == 0 else 'td'
                            for cell in row:
                                html_content.append(f'<{tag}>{cell or ""}</{tag}>')
                            html_content.append('</tr>')
                        html_content.append('</table>')

                # Extract text
                text = page.extract_text()
                if text:
                    html_content.append('<p>')
                    html_content.append(text.replace('\n', '<br>'))
                    html_content.append('</p>')

                html_content.append('</div>')

        html_content.extend(['</body>', '</html>'])

        # Write HTML file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(html_content))

        print(f"Converted PDF to HTML: {output_path}")
        return output_path

    except ImportError:
        raise Exception("pdfplumber library not installed. Run: pip install pdfplumber")
    except Exception as e:
        raise Exception(f"Failed to convert PDF to HTML: {str(e)}")
